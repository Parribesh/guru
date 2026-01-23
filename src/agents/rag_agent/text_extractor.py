from __future__ import annotations

from pathlib import Path
from typing import Optional


class UnsupportedDocumentTypeError(ValueError):
    pass


def extract_text_from_path(path: str | Path) -> str:
    """
    Extract plain text from common document formats.

    Supported (best-effort):
    - .txt/.md/.json/.csv/... (generic text)
    - .pdf (via `pypdf` if installed)
    - .docx (via `python-docx` if installed)
    - .pptx (via `python-pptx` if installed)
    - .rtf (via `striprtf` if installed)
    - .html/.htm (via `beautifulsoup4` if installed)

    Notes:
    - This intentionally avoids shelling out to system tools.
    - For unsupported formats, it raises UnsupportedDocumentTypeError.
    """

    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(str(p))

    suffix = p.suffix.lower().lstrip(".")
    if suffix in {"", "txt", "md", "markdown", "csv", "tsv", "json", "xml", "yaml", "yml", "py", "js", "ts"}:
        return _read_text_best_effort(p)

    if suffix == "pdf":
        return _extract_pdf(p)

    if suffix == "docx":
        return _extract_docx(p)

    if suffix == "pptx":
        return _extract_pptx(p)

    if suffix == "rtf":
        return _extract_rtf(p)

    if suffix in {"html", "htm"}:
        return _extract_html(p)

    # Fallback: try reading as text; if it's garbage, ask caller to handle.
    text = _read_text_best_effort(p, raise_on_decode_error=False)
    if text.strip():
        return text
    raise UnsupportedDocumentTypeError(f"Unsupported document type: .{suffix}")


def _read_text_best_effort(p: Path, *, raise_on_decode_error: bool = True) -> str:
    raw = p.read_bytes()
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        if raise_on_decode_error:
            # best-effort decode for unknown encodings
            return raw.decode("utf-8", errors="ignore")
        return ""


def _extract_pdf(p: Path) -> str:
    try:
        from pypdf import PdfReader  # type: ignore
    except ModuleNotFoundError as e:
        raise UnsupportedDocumentTypeError("PDF extraction requires `pypdf`") from e

    reader = PdfReader(str(p))
    parts: list[str] = []
    for page in reader.pages:
        txt = page.extract_text() or ""
        if txt.strip():
            parts.append(txt)
    return "\n\n".join(parts).strip()


def _extract_docx(p: Path) -> str:
    try:
        from docx import Document  # type: ignore
    except ModuleNotFoundError as e:
        raise UnsupportedDocumentTypeError("DOCX extraction requires `python-docx`") from e

    doc = Document(str(p))
    parts: list[str] = []
    for para in doc.paragraphs:
        t = (para.text or "").strip()
        if t:
            parts.append(t)
    return "\n".join(parts).strip()


def _extract_pptx(p: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ModuleNotFoundError as e:
        raise UnsupportedDocumentTypeError("PPTX extraction requires `python-pptx`") from e

    prs = Presentation(str(p))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            text: Optional[str] = getattr(shape, "text", None)
            if text and text.strip():
                parts.append(text.strip())
    return "\n\n".join(parts).strip()


def _extract_rtf(p: Path) -> str:
    try:
        from striprtf.striprtf import rtf_to_text  # type: ignore
    except ModuleNotFoundError as e:
        raise UnsupportedDocumentTypeError("RTF extraction requires `striprtf`") from e

    raw = p.read_text(encoding="utf-8", errors="ignore")
    return (rtf_to_text(raw) or "").strip()


def _extract_html(p: Path) -> str:
    try:
        from bs4 import BeautifulSoup  # type: ignore
    except ModuleNotFoundError as e:
        raise UnsupportedDocumentTypeError("HTML extraction requires `beautifulsoup4`") from e

    raw = p.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(raw, "html.parser")
    return soup.get_text(separator="\n").strip()


